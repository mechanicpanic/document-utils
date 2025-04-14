import os
import sys
import json
import argparse
import time
from pptx import Presentation
from google import genai
from tqdm import tqdm
import tempfile
import subprocess


def pdf_to_pptx(pdf_path, output_path=None):
    """
    Convert PDF to PPTX using pdf2pptx library.
    If pdf2pptx is not installed, it will be installed automatically.
    """
    try:
        from pdf2pptx import convert_pdf2pptx
    except ImportError:
        print("Installing pdf2pptx library...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "pdf2pptx"])
        from pdf2pptx import convert_pdf2pptx

    if output_path is None:
        base_name = os.path.splitext(pdf_path)[0]
        output_path = f"{base_name}.pptx"

    print(f"Converting PDF to PPTX: {pdf_path} -> {output_path}")
    convert_pdf2pptx(pdf_path, output_path)
    return output_path


def extract_text_from_pptx(pptx_path):
    """Extract text content from a PPTX file."""
    presentation = Presentation(pptx_path)
    slides_content = []

    for slide_number, slide in enumerate(presentation.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(
                    {"text": shape.text.strip(), "shape_id": shape.shape_id}
                )

        if slide_texts:
            slides_content.append(
                {"slide_number": slide_number, "content": slide_texts}
            )

    return slides_content


def translate_text_with_gemini(
    text, target_language, client, retry_count=3, wait_time=1.0
):
    """
    Translate text using the Google Gemini API with rate limiting and retries.

    Args:
        text: Text to translate
        target_language: Target language
        client: Gemini API client
        retry_count: Number of retries on failure
        wait_time: Base wait time between API calls (seconds)
    """
    prompt = f"""You are a professional translator. Translate the following text to {target_language}. 
Preserve formatting like line breaks, bullet points, and paragraph structure.
Return ONLY the translated text without any explanations or notes.

Text to translate:
{text}"""

    attempts = 0
    while attempts < retry_count:
        try:
            # Add wait time before API call to respect rate limits
            time.sleep(wait_time)

            response = client.models.generate_content(
                model="gemini-2.0-flash-lite", contents=prompt
            )
            return response.text

        except Exception as e:
            attempts += 1
            error_type = str(e)

            # Handle rate limit errors with exponential backoff
            if "rate limit" in error_type.lower() or "quota" in error_type.lower():
                backoff_time = wait_time * (2**attempts)
                print(
                    f"Rate limit exceeded. Waiting {backoff_time:.2f} seconds before retry..."
                )
                time.sleep(backoff_time)
            else:
                print(
                    f"Error during translation (attempt {attempts}/{retry_count}): {e}"
                )

            # If we've exhausted retries, return original text
            if attempts >= retry_count:
                print(
                    f"Failed to translate after {retry_count} attempts. Returning original text."
                )
                return text

    return text  # Return original text if translation fails


def replace_text_in_pptx(pptx_path, translated_content, output_path):
    """Replace original text with translated text in the PPTX file."""
    presentation = Presentation(pptx_path)

    for slide_info in translated_content:
        slide_number = slide_info["slide_number"]
        if slide_number <= len(presentation.slides):
            slide = presentation.slides[slide_number - 1]

            shape_id_map = {
                shape.shape_id: shape
                for shape in slide.shapes
                if hasattr(shape, "text")
            }

            for content_item in slide_info["content"]:
                shape_id = content_item["shape_id"]
                translated_text = content_item["translated_text"]

                if shape_id in shape_id_map:
                    shape = shape_id_map[shape_id]
                    if hasattr(shape, "text"):
                        shape.text = translated_text

    presentation.save(output_path)


def process_batch(text_batch, target_language, client, retry_count=3, wait_time=1.0):
    """
    Process a batch of text blocks in parallel to optimize API usage.

    This batching approach can be further optimized with real parallelism
    using threading or asyncio if needed.
    """
    results = []
    for text in text_batch:
        translated = translate_text_with_gemini(
            text, target_language, client, retry_count, wait_time
        )
        results.append(translated)
    return results


def main():
    parser = argparse.ArgumentParser(
        description="Convert PDF to PPTX and translate using Google Gemini"
    )
    parser.add_argument("input_file", help="Path to the input file (PDF or PPTX)")
    parser.add_argument("target_language", help="Target language for translation")
    parser.add_argument("--output_file", help="Path to save the translated file")
    parser.add_argument("--api_key", help="Google API key")
    parser.add_argument(
        "--skip_conversion",
        action="store_true",
        help="Skip PDF to PPTX conversion (if input is already PPTX)",
    )
    parser.add_argument(
        "--wait_time",
        type=float,
        default=1.0,
        help="Wait time between API calls in seconds (default: 1.0)",
    )
    parser.add_argument(
        "--retry_count",
        type=int,
        default=3,
        help="Number of retries on API failure (default: 3)",
    )
    parser.add_argument(
        "--batch_size",
        type=int,
        default=1,
        help="Number of text blocks to process in each batch (default: 1)",
    )
    parser.add_argument(
        "--save_progress",
        action="store_true",
        help="Save translation progress to resume later if interrupted",
    )

    args = parser.parse_args()

    # Use environment variable if API key not provided as argument
    api_key = args.api_key or os.environ.get("GOOGLE_API_KEY")
    if not api_key:
        raise ValueError(
            "Google API key not provided. Use --api_key or set GOOGLE_API_KEY environment variable."
        )

    # Initialize Gemini client
    client = genai.Client(api_key=api_key)

    # Determine file type and convert if necessary
    file_ext = os.path.splitext(args.input_file)[1].lower()
    if file_ext == ".pdf" and not args.skip_conversion:
        # Convert PDF to PPTX
        temp_pptx = pdf_to_pptx(args.input_file)
        pptx_path = temp_pptx
    elif file_ext == ".pptx":
        pptx_path = args.input_file
    else:
        raise ValueError(
            f"Unsupported file format: {file_ext}. Please provide a PDF or PPTX file."
        )

    # Set default output file if not specified
    if not args.output_file:
        base_name = os.path.splitext(args.input_file)[0]
        args.output_file = f"{base_name}_{args.target_language}.pptx"

    print(f"Extracting text from {pptx_path}...")
    slides_content = extract_text_from_pptx(pptx_path)

    print(f"Translating content to {args.target_language}...")
    total_text_blocks = sum(len(slide_info["content"]) for slide_info in slides_content)
    print(f"Total text blocks to translate: {total_text_blocks}")

    # Create progress file path for saving/resuming
    progress_file = None
    if args.save_progress:
        progress_file = f"{os.path.splitext(args.output_file)[0]}_progress.json"

        # Check if we can resume from existing progress
        if os.path.exists(progress_file):
            try:
                with open(progress_file, "r", encoding="utf-8") as f:
                    progress_data = json.load(f)

                # Apply saved translations to our slides_content
                applied_count = 0
                for slide_info in slides_content:
                    slide_num = slide_info["slide_number"]
                    if str(slide_num) in progress_data:
                        for content_item in slide_info["content"]:
                            shape_id = content_item["shape_id"]
                            if str(shape_id) in progress_data[str(slide_num)]:
                                content_item["translated_text"] = progress_data[
                                    str(slide_num)
                                ][str(shape_id)]
                                applied_count += 1

                if applied_count > 0:
                    print(
                        f"Resumed translation progress: {applied_count}/{total_text_blocks} blocks already translated"
                    )
            except Exception as e:
                print(f"Error loading progress file: {e}")

    # Collect all text blocks that need translation
    to_translate = []
    for slide_info in slides_content:
        for content_item in slide_info["content"]:
            if "translated_text" not in content_item:
                to_translate.append(
                    (
                        slide_info["slide_number"],
                        content_item["shape_id"],
                        content_item["text"],
                    )
                )

    # Process in batches
    batch_size = max(1, min(args.batch_size, len(to_translate)))
    with tqdm(total=len(to_translate), desc="Translating") as pbar:
        for i in range(0, len(to_translate), batch_size):
            batch = to_translate[i : i + batch_size]

            # Extract just the text for translation
            text_batch = [item[2] for item in batch]

            # Translate the batch
            translated_batch = []
            for text in text_batch:
                translated = translate_text_with_gemini(
                    text,
                    args.target_language,
                    client,
                    retry_count=args.retry_count,
                    wait_time=args.wait_time,
                )
                translated_batch.append(translated)
                pbar.update(1)

            # Apply translations back to the content
            for j, (slide_num, shape_id, _) in enumerate(batch):
                # Find the correct slide and content item
                for slide_info in slides_content:
                    if slide_info["slide_number"] == slide_num:
                        for content_item in slide_info["content"]:
                            if content_item["shape_id"] == shape_id:
                                content_item["translated_text"] = translated_batch[j]
                                break

            # Save progress if requested
            if progress_file:
                progress_data = {}
                for slide_info in slides_content:
                    slide_num = str(slide_info["slide_number"])
                    progress_data[slide_num] = {}

                    for content_item in slide_info["content"]:
                        if "translated_text" in content_item:
                            shape_id = str(content_item["shape_id"])
                            progress_data[slide_num][shape_id] = content_item[
                                "translated_text"
                            ]

                with open(progress_file, "w", encoding="utf-8") as f:
                    json.dump(progress_data, f, ensure_ascii=False, indent=2)

    print(f"Replacing text in the presentation...")
    replace_text_in_pptx(pptx_path, slides_content, args.output_file)

    # Clean up temporary file if needed
    if (
        file_ext == ".pdf"
        and pptx_path != args.input_file
        and os.path.exists(pptx_path)
    ):
        if pptx_path != args.output_file:  # Don't delete if it's the output file
            os.remove(pptx_path)

    print(f"Translation completed. Output saved to {args.output_file}")


if __name__ == "__main__":
    main()
